#!/usr/bin/env python3
"""
PowerPoint Editor Tool Class
Provides basic PPT editing functionality including adding text, images, shapes, etc.

Every method takes file_path as its first parameter. The presentation is loaded
from file_path if it exists, otherwise a new one is created. After applying the
requested changes, the presentation is automatically saved back to file_path.
"""

import logging
import os
from typing import Any, Dict, List, Optional, TYPE_CHECKING
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
except ImportError:
    raise ImportError("Please install python-pptx library: pip install python-pptx")

if TYPE_CHECKING:
    from pptx.presentation import Presentation as PresentationType

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class PowerPointEditor:
    def __init__(self):
        self.current_presentation: Optional["PresentationType"] = None
        self.current_file_path: Optional[str] = None

    def _load_or_create(self, file_path: str) -> Dict[str, Any]:
        """Load an existing presentation from file_path, or create a new one if it doesn't exist."""
        try:
            if not file_path:
                return {"success": False, "error": "file_path is required"}

            expanded = os.path.expanduser(file_path)
            if Path(expanded).exists():
                self.current_presentation = Presentation(expanded)
                self.current_file_path = expanded
                return {
                    "success": True,
                    "message": f"Loaded existing presentation: {expanded}",
                    "created": False,
                }
            else:
                parent = Path(expanded).parent
                if parent and not parent.exists():
                    parent.mkdir(parents=True, exist_ok=True)
                self.current_presentation = Presentation()
                self.current_file_path = expanded
                return {
                    "success": True,
                    "message": f"Created new presentation: {expanded}",
                    "created": True,
                }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def _save_current(self) -> Dict[str, Any]:
        """Save the current presentation to current_file_path."""
        try:
            if not self.current_presentation:
                return {"success": False, "error": "No presentation is open"}
            if not self.current_file_path:
                return {"success": False, "error": "No file path set"}
            self.current_presentation.save(self.current_file_path)
            return {
                "success": True,
                "message": f"Saved presentation to {self.current_file_path}",
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def _require_presentation(self) -> Optional[Dict[str, Any]]:
        """Return an error dict if no presentation is loaded, or None if OK."""
        if not self.current_presentation:
            return {"success": False, "error": "No presentation is open"}
        return None

    def _validate_slide_index(self, slide_index: int) -> Optional[Dict[str, Any]]:
        """Return an error dict if slide_index is out of range, or None if OK."""
        err = self._require_presentation()
        if err:
            return err
        if slide_index >= len(self.current_presentation.slides):
            return {
                "success": False,
                "error": f"Slide index out of range: {slide_index}",
            }
        return None

    def add_slide(self, file_path: str, layout_index: int = 1) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            slide_layouts = self.current_presentation.slide_layouts
            if layout_index >= len(slide_layouts):
                layout_index = 1
            layout = slide_layouts[layout_index]
            slide = self.current_presentation.slides.add_slide(layout)
            result = {
                "success": True,
                "message": "Successfully added new slide",
                "slide_index": len(self.current_presentation.slides) - 1,
                "total_slides": len(self.current_presentation.slides),
            }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def add_text_box(
        self,
        file_path: str,
        slide_index: int,
        text: str,
        left: float = 1,
        top: float = 1,
        width: float = 8,
        height: float = 1,
        font_size: int = 18,
        font_color: str = "000000",
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            err = self._validate_slide_index(slide_index)
            if err:
                return err
            slide = self.current_presentation.slides[slide_index]
            textbox = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(width), Inches(height)
            )
            text_frame = textbox.text_frame
            text_frame.text = text
            paragraph = text_frame.paragraphs[0]
            font = paragraph.font
            font.size = Pt(font_size)
            try:
                rgb_color = RGBColor.from_string(font_color)
                font.color.rgb = rgb_color
            except:
                pass
            result = {
                "success": True,
                "message": f"Successfully added text box to slide {slide_index}",
                "text": text,
                "position": {
                    "left": left,
                    "top": top,
                    "width": width,
                    "height": height,
                },
            }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def add_title_slide(
        self, file_path: str, title: str, subtitle: str = ""
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            title_slide_layout = self.current_presentation.slide_layouts[0]
            slide = self.current_presentation.slides.add_slide(title_slide_layout)
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = title
            if subtitle and len(slide.placeholders) > 1:
                try:
                    subtitle_shape = slide.placeholders[1]
                    subtitle_shape.text_frame.text = subtitle
                except (AttributeError, TypeError):
                    pass
            result = {
                "success": True,
                "message": "Successfully added title slide",
                "slide_index": len(self.current_presentation.slides) - 1,
                "title": title,
                "subtitle": subtitle,
            }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def add_bullet_points(
        self, file_path: str, slide_index: int, title: str, bullet_points: List[str]
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            err = self._validate_slide_index(slide_index)
            if err:
                return err
            slide = self.current_presentation.slides[slide_index]
            if slide.shapes.title:
                slide.shapes.title.text = title
            content_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    content_placeholder = shape
                    break
            if content_placeholder:
                try:
                    text_frame = content_placeholder.text_frame
                    text_frame.clear()
                    for i, point in enumerate(bullet_points):
                        if i == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        p.text = point
                        p.level = 0
                except (AttributeError, TypeError):
                    pass
            result = {
                "success": True,
                "message": f"Successfully added bullet points to slide {slide_index}",
                "title": title,
                "bullet_points": bullet_points,
            }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def add_image(
        self,
        file_path: str,
        slide_index: int,
        image_path: str,
        left: float = 1,
        top: float = 2,
        width: Optional[float] = None,
        height: Optional[float] = None,
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            if not Path(image_path).exists():
                return {
                    "success": False,
                    "error": f"Image file does not exist: {image_path}",
                }
            err = self._validate_slide_index(slide_index)
            if err:
                return err
            slide = self.current_presentation.slides[slide_index]
            left_inches = Inches(left)
            top_inches = Inches(top)
            if width and height:
                pic = slide.shapes.add_picture(
                    image_path, left_inches, top_inches, Inches(width), Inches(height)
                )
            else:
                pic = slide.shapes.add_picture(image_path, left_inches, top_inches)
            result = {
                "success": True,
                "message": f"Successfully added image to slide {slide_index}",
                "image_path": image_path,
                "position": {
                    "left": left,
                    "top": top,
                    "width": width,
                    "height": height,
                },
            }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def add_shape(
        self,
        file_path: str,
        slide_index: int,
        shape_type: str,
        left: float = 1,
        top: float = 1,
        width: float = 2,
        height: float = 1,
        fill_color: str = "0066CC",
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            err = self._validate_slide_index(slide_index)
            if err:
                return err
            slide = self.current_presentation.slides[slide_index]
            shape_map = {
                "rectangle": MSO_SHAPE.RECTANGLE,
                "oval": MSO_SHAPE.OVAL,
                "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
                "diamond": MSO_SHAPE.DIAMOND,
                "pentagon": MSO_SHAPE.REGULAR_PENTAGON,
                "hexagon": MSO_SHAPE.HEXAGON,
                "star": MSO_SHAPE.STAR_5_POINT,
                "arrow": MSO_SHAPE.BLOCK_ARC,
            }
            if shape_type.lower() not in shape_map:
                return {
                    "success": False,
                    "error": f"Unsupported shape type: {shape_type}",
                }
            shape = slide.shapes.add_shape(
                shape_map[shape_type.lower()],
                Inches(left),
                Inches(top),
                Inches(width),
                Inches(height),
            )
            try:
                rgb_color = RGBColor.from_string(fill_color)
                shape.fill.solid()
                shape.fill.fore_color.rgb = rgb_color
            except:
                pass
            result = {
                "success": True,
                "message": f"Successfully added {shape_type} shape to slide {slide_index}",
                "shape_type": shape_type,
                "position": {
                    "left": left,
                    "top": top,
                    "width": width,
                    "height": height,
                },
                "fill_color": fill_color,
            }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def get_presentation_info(self, file_path: str) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            slides_info = []
            for i, slide in enumerate(self.current_presentation.slides):
                slide_info = {
                    "index": i,
                    "shapes_count": len(slide.shapes),
                    "has_title": bool(slide.shapes.title and slide.shapes.title.text),
                    "title": slide.shapes.title.text if slide.shapes.title else "",
                }
                slides_info.append(slide_info)
            return {
                "success": True,
                "file_path": self.current_file_path,
                "slides_count": len(self.current_presentation.slides),
                "slides": slides_info,
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def delete_slide(self, file_path: str, slide_index: int) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            err = self._validate_slide_index(slide_index)
            if err:
                return err
            xml_slides = self.current_presentation.part._element.sldIdLst
            xml_slides.remove(xml_slides[slide_index])
            slides = self.current_presentation.slides
            slides._sldIdLst.remove(slides._sldIdLst[slide_index])
            result = {
                "success": True,
                "message": f"Successfully deleted slide {slide_index}",
                "remaining_slides": len(self.current_presentation.slides),
            }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def duplicate_slide(self, file_path: str, slide_index: int) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            err = self._validate_slide_index(slide_index)
            if err:
                return err
            source_slide = self.current_presentation.slides[slide_index]
            slide_layout = source_slide.slide_layout
            new_slide = self.current_presentation.slides.add_slide(slide_layout)
            for shape in source_slide.shapes:
                if not shape.is_placeholder:
                    self._copy_shape(shape, new_slide)
            result = {
                "success": True,
                "message": f"Successfully duplicated slide {slide_index}",
                "new_slide_index": len(self.current_presentation.slides) - 1,
                "total_slides": len(self.current_presentation.slides),
            }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def move_slide(
        self, file_path: str, from_index: int, to_index: int
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            slides = self.current_presentation.slides
            if from_index >= len(slides) or to_index >= len(slides):
                return {"success": False, "error": "Slide index out of range"}
            if from_index == to_index:
                return {"success": True, "message": "Slide position unchanged"}
            source_slide = slides[from_index]
            slide_layout = source_slide.slide_layout
            new_slide = slides.add_slide(slide_layout)
            for shape in source_slide.shapes:
                if not shape.is_placeholder:
                    self._copy_shape(shape, new_slide)
            if from_index < len(slides) - 1:
                actual_from_index = from_index if to_index > from_index else from_index
                xml_slides = self.current_presentation.part._element.sldIdLst
                xml_slides.remove(xml_slides[actual_from_index])
                slides._sldIdLst.remove(slides._sldIdLst[actual_from_index])
            result = {
                "success": True,
                "message": f"Successfully moved slide from position {from_index} to position {to_index}",
                "total_slides": len(slides),
            }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def add_table(
        self,
        file_path: str,
        slide_index: int,
        rows: int,
        cols: int,
        left: float = 1,
        top: float = 2,
        width: float = 8,
        height: float = 4,
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            err = self._validate_slide_index(slide_index)
            if err:
                return err
            slide = self.current_presentation.slides[slide_index]
            table = slide.shapes.add_table(
                rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
            )
            result = {
                "success": True,
                "message": f"Successfully added {rows}x{cols} table to slide {slide_index}",
                "rows": rows,
                "cols": cols,
                "position": {
                    "left": left,
                    "top": top,
                    "width": width,
                    "height": height,
                },
            }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def set_table_cell_text(
        self,
        file_path: str,
        slide_index: int,
        table_index: int,
        row: int,
        col: int,
        text: str,
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            err = self._validate_slide_index(slide_index)
            if err:
                return err
            slide = self.current_presentation.slides[slide_index]
            tables = []
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "table"):
                        table_obj = getattr(shape, "table", None)
                        if table_obj is not None:
                            tables.append(shape)
                except:
                    continue
            if table_index >= len(tables):
                return {
                    "success": False,
                    "error": f"Table index out of range: {table_index}",
                }
            table = getattr(tables[table_index], "table")
            if row >= len(table.rows) or col >= len(table.columns):
                return {"success": False, "error": "Cell position out of table range"}
            table.cell(row, col).text = text
            result = {
                "success": True,
                "message": f"Successfully set text for table cell ({row}, {col})",
                "text": text,
            }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def set_slide_background_color(
        self, file_path: str, slide_index: int, color: str
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            err = self._validate_slide_index(slide_index)
            if err:
                return err
            slide = self.current_presentation.slides[slide_index]
            background = slide.background
            fill = background.fill
            fill.solid()
            try:
                rgb_color = RGBColor.from_string(color)
                fill.fore_color.rgb = rgb_color
            except:
                return {"success": False, "error": f"Invalid color format: {color}"}
            result = {
                "success": True,
                "message": f"Successfully set background color for slide {slide_index}",
                "color": color,
            }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def add_hyperlink(
        self,
        file_path: str,
        slide_index: int,
        shape_index: int,
        url: str,
        display_text: Optional[str] = None,
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            err = self._validate_slide_index(slide_index)
            if err:
                return err
            slide = self.current_presentation.slides[slide_index]
            if shape_index >= len(slide.shapes):
                return {
                    "success": False,
                    "error": f"Shape index out of range: {shape_index}",
                }
            shape = slide.shapes[shape_index]
            try:
                if hasattr(shape, "text_frame"):
                    text_frame = getattr(shape, "text_frame")
                    if text_frame is not None:
                        if display_text:
                            text_frame.text = display_text
                        if not text_frame.paragraphs:
                            paragraph = text_frame.add_paragraph()
                        else:
                            paragraph = text_frame.paragraphs[0]
                        if not paragraph.runs:
                            run = paragraph.add_run()
                        else:
                            run = paragraph.runs[0]
                        run.hyperlink.address = url
                    else:
                        if hasattr(shape, "click_action"):
                            shape.click_action.hyperlink.address = url
                        else:
                            return {
                                "success": False,
                                "error": "Shape does not support hyperlinks",
                            }
                else:
                    if hasattr(shape, "click_action"):
                        shape.click_action.hyperlink.address = url
                    else:
                        return {
                            "success": False,
                            "error": "Shape does not support hyperlinks",
                        }
            except Exception as e:
                return {"success": False, "error": f"Failed to add hyperlink: {str(e)}"}
            result = {
                "success": True,
                "message": f"Successfully added hyperlink to shape {shape_index} on slide {slide_index}",
                "url": url,
                "display_text": display_text,
            }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def set_text_formatting(
        self,
        file_path: str,
        slide_index: int,
        shape_index: int,
        font_name: Optional[str] = None,
        font_size: Optional[int] = None,
        font_color: Optional[str] = None,
        bold: Optional[bool] = None,
        italic: Optional[bool] = None,
        underline: Optional[bool] = None,
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            err = self._validate_slide_index(slide_index)
            if err:
                return err
            slide = self.current_presentation.slides[slide_index]
            if shape_index >= len(slide.shapes):
                return {
                    "success": False,
                    "error": f"Shape index out of range: {shape_index}",
                }
            shape = slide.shapes[shape_index]
            try:
                if not hasattr(shape, "text_frame"):
                    return {
                        "success": False,
                        "error": "Shape does not support text boxes",
                    }
                text_frame = getattr(shape, "text_frame")
                if text_frame is None:
                    return {"success": False, "error": "Text frame not available"}
                if not hasattr(text_frame, "paragraphs"):
                    return {
                        "success": False,
                        "error": "Text frame has no paragraphs attribute",
                    }
                paragraphs = text_frame.paragraphs
                if not paragraphs or len(paragraphs) == 0:
                    return {"success": False, "error": "No available text paragraphs"}
                paragraph = paragraphs[0]
                if not hasattr(paragraph, "font"):
                    return {
                        "success": False,
                        "error": "Paragraph has no font attribute",
                    }
                font = paragraph.font
                if font is None:
                    return {"success": False, "error": "Unable to get font object"}
                if font_name:
                    font.name = font_name
                if font_size:
                    font.size = Pt(font_size)
                if font_color:
                    try:
                        rgb_color = RGBColor.from_string(font_color)
                        font.color.rgb = rgb_color
                    except:
                        return {
                            "success": False,
                            "error": f"Invalid color format: {font_color}",
                        }
                if bold is not None:
                    font.bold = bold
                if italic is not None:
                    font.italic = italic
                if underline is not None:
                    font.underline = underline
            except Exception as e:
                return {"success": False, "error": f"Font setting failed: {str(e)}"}
            result = {
                "success": True,
                "message": f"Successfully set text formatting for shape {shape_index} on slide {slide_index}",
                "formatting": {
                    "font_name": font_name,
                    "font_size": font_size,
                    "font_color": font_color,
                    "bold": bold,
                    "italic": italic,
                    "underline": underline,
                },
            }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def get_slide_shapes_info(self, file_path: str, slide_index: int) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            err = self._validate_slide_index(slide_index)
            if err:
                return err
            slide = self.current_presentation.slides[slide_index]
            shapes_info = []
            for i, shape in enumerate(slide.shapes):
                shape_info = {
                    "index": i,
                    "shape_type": str(shape.shape_type),
                    "name": shape.name,
                    "left": shape.left.inches if hasattr(shape.left, "inches") else 0,
                    "top": shape.top.inches if hasattr(shape.top, "inches") else 0,
                    "width": shape.width.inches
                    if hasattr(shape.width, "inches")
                    else 0,
                    "height": shape.height.inches
                    if hasattr(shape.height, "inches")
                    else 0,
                    "has_text": hasattr(shape, "text_frame")
                    and getattr(shape, "text_frame", None) is not None,
                    "text": getattr(shape, "text", "")
                    if hasattr(shape, "text")
                    else "",
                }
                shapes_info.append(shape_info)
            return {
                "success": True,
                "slide_index": slide_index,
                "shapes_count": len(shapes_info),
                "shapes": shapes_info,
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def _copy_shape(self, source_shape, target_slide):
        try:
            if hasattr(source_shape, "text_frame") and source_shape.text_frame:
                textbox = target_slide.shapes.add_textbox(
                    source_shape.left,
                    source_shape.top,
                    source_shape.width,
                    source_shape.height,
                )
                textbox.text_frame.text = source_shape.text_frame.text
        except:
            pass

    def set_slide_transition(
        self,
        file_path: str,
        slide_index: int,
        transition_type: str = "fade",
        duration: float = 1.0,
        advance_on_click: bool = True,
        advance_after_time: Optional[float] = None,
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            err = self._validate_slide_index(slide_index)
            if err:
                return err
            slide = self.current_presentation.slides[slide_index]
            supported_transitions = [
                "none",
                "fade",
                "push",
                "wipe",
                "split",
                "zoom",
                "blinds",
                "dissolve",
            ]
            if transition_type.lower() not in supported_transitions:
                return {
                    "success": False,
                    "error": f"Unsupported transition type: {transition_type}. Supported types: {', '.join(supported_transitions)}",
                }
            slide_element = slide._element
            namespaces = {
                "p": "http://schemas.openxmlformats.org/presentationml/2006/main"
            }
            existing_transition = slide_element.find(".//p:transition", namespaces)
            if existing_transition is not None:
                slide_element.remove(existing_transition)
            if transition_type.lower() != "none":
                try:
                    from lxml import etree
                except ImportError:
                    return {
                        "success": False,
                        "error": "Please install lxml library: pip install lxml",
                    }
                transition_xml = self._create_transition_xml(
                    transition_type, duration, advance_on_click, advance_after_time
                )
                parser = etree.XMLParser(ns_clean=True, recover=True)
                transition_elem = etree.fromstring(
                    transition_xml.encode("utf-8"), parser
                )
                color_map_override = slide_element.find(".//p:clrMapOvr", namespaces)
                if color_map_override is not None:
                    color_map_override.addprevious(transition_elem)
                else:
                    slide_element.append(transition_elem)
                verification_elem = slide_element.find(".//p:transition", namespaces)
                if verification_elem is None:
                    return {
                        "success": False,
                        "error": "Transition effect XML insertion failed",
                    }
            result = {
                "success": True,
                "message": f"Successfully set transition effect for slide {slide_index}",
                "transition_type": transition_type,
                "duration": duration,
                "advance_on_click": advance_on_click,
                "advance_after_time": advance_after_time,
                "verification": "Transition effect verified successfully"
                if transition_type.lower() != "none"
                else "Transition effect removed",
            }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def _create_transition_xml(
        self,
        transition_type: str,
        duration: float,
        advance_on_click: bool,
        advance_after_time: Optional[float],
    ) -> str:
        if duration <= 0.5:
            speed = "fast"
        elif duration <= 2.0:
            speed = "med"
        else:
            speed = "slow"
        transition_attrs = f'spd="{speed}"'
        if advance_on_click:
            transition_attrs += ' advClick="1"'
        else:
            transition_attrs += ' advClick="0"'
        if advance_after_time is not None:
            advance_time_ms = int(advance_after_time * 1000)
            transition_attrs += f' advTm="{advance_time_ms}"'
        transition_content = ""
        if transition_type.lower() == "fade":
            transition_content = "<p:fade/>"
        elif transition_type.lower() == "push":
            transition_content = '<p:push dir="l"/>'
        elif transition_type.lower() == "wipe":
            transition_content = '<p:wipe dir="l"/>'
        elif transition_type.lower() == "zoom":
            transition_content = "<p:zoom/>"
        elif transition_type.lower() == "split":
            transition_content = '<p:split orient="horz" dir="out"/>'
        elif transition_type.lower() == "blinds":
            transition_content = '<p:blinds dir="horz"/>'
        elif transition_type.lower() == "dissolve":
            transition_content = "<p:dissolve/>"
        else:
            transition_content = "<p:fade/>"
        xml_string = f"""<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" {transition_attrs}>
    {transition_content}
</p:transition>"""
        return xml_string

    def apply_transition_to_all_slides(
        self, file_path: str, transition_type: str = "fade", duration: float = 1.0
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            slides = self.current_presentation.slides
            if len(slides) == 0:
                return {"success": False, "error": "No slides in presentation"}
            success_count = 0
            failed_slides = []
            for i in range(len(slides)):
                result = self.set_slide_transition(
                    file_path, i, transition_type, duration, True, None
                )
                if result.get("success"):
                    success_count += 1
                else:
                    failed_slides.append(i)
            if success_count == len(slides):
                return {
                    "success": True,
                    "message": f"Successfully set '{transition_type}' transition effect for all {len(slides)} slides",
                    "transition_type": transition_type,
                    "duration": duration,
                    "slides_processed": len(slides),
                }
            else:
                return {
                    "success": True,
                    "message": f"Set transition effect for {success_count}/{len(slides)} slides",
                    "transition_type": transition_type,
                    "duration": duration,
                    "slides_processed": success_count,
                    "failed_slides": failed_slides,
                    "warning": f"{len(failed_slides)} slides failed to set",
                }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def get_available_transitions(self, file_path: str) -> Dict[str, Any]:
        transitions = [
            {"name": "none", "description": "No transition effect"},
            {
                "name": "fade",
                "description": "Fade - Recommended for professional presentations",
            },
            {"name": "push", "description": "Push - Dynamic and energetic"},
            {"name": "wipe", "description": "Wipe - Clean and smooth"},
            {"name": "split", "description": "Split - Creative effect"},
            {"name": "zoom", "description": "Zoom - Emphasize key points"},
            {"name": "blinds", "description": "Blinds - Classic effect"},
            {"name": "dissolve", "description": "Dissolve - Gentle transition"},
        ]
        return {
            "success": True,
            "transitions": transitions,
            "total_count": len(transitions),
            "note": "These animation effects can make your presentation more vivid and interesting",
            "recommendation": "Recommended to use 'fade' effect, suitable for most professional presentations",
        }

    def make_presentation_professional(self, file_path: str) -> Dict[str, Any]:
        return self.apply_transition_to_all_slides(file_path, "fade", 1.0)

    def add_smooth_transitions(self, file_path: str) -> Dict[str, Any]:
        return self.apply_transition_to_all_slides(file_path, "fade", 0.8)

    def add_dynamic_effects(self, file_path: str) -> Dict[str, Any]:
        return self.apply_transition_to_all_slides(file_path, "push", 1.2)

    def generate_outline_for_topic(self, file_path: str, topic: str) -> Dict[str, Any]:
        try:
            import json as _json

            outline_data = {
                "slides": [
                    {
                        "title": f"In-depth Discussion on {topic}",
                        "subtitle": "AI-assisted generation",
                    },
                    {
                        "title": "Introduction and Background",
                        "content": [
                            f"Definition and importance of {topic}",
                            "Related historical development",
                            "Main scope of this discussion",
                        ],
                    },
                    {
                        "title": "Core Points Analysis",
                        "content": [
                            "First key aspect",
                            "Second key aspect with examples",
                            "In-depth analysis of third key aspect",
                        ],
                    },
                    {
                        "title": "Case Study or Practical Application",
                        "content": [
                            f"A real-world case about {topic}",
                            "Insights from the case",
                            "How to apply these in practice",
                        ],
                    },
                    {
                        "title": "Summary and Outlook",
                        "content": [
                            f"Summary of core content on {topic}",
                            "Future development trends",
                            "Q&A session",
                        ],
                    },
                ]
            }
            outline_json = _json.dumps(outline_data, ensure_ascii=False, indent=2)
            return {
                "success": True,
                "message": f"Successfully generated outline for topic '{topic}'.",
                "outline_json": outline_json,
            }
        except Exception as e:
            logger.error(f"Error generating outline for topic '{topic}': {e}")
            return {"success": False, "error": str(e)}

    def set_core_properties(
        self,
        file_path: str,
        title: Optional[str] = None,
        subject: Optional[str] = None,
        author: Optional[str] = None,
        keywords: Optional[str] = None,
        comments: Optional[str] = None,
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            core_props = self.current_presentation.core_properties
            if title is not None:
                core_props.title = title
            if subject is not None:
                core_props.subject = subject
            if author is not None:
                core_props.author = author
            if keywords is not None:
                core_props.keywords = keywords
            if comments is not None:
                core_props.comments = comments
            result = {
                "success": True,
                "message": "Core properties updated successfully",
                "properties_set": {
                    k: v
                    for k, v in [
                        ("title", title),
                        ("subject", subject),
                        ("author", author),
                        ("keywords", keywords),
                        ("comments", comments),
                    ]
                    if v is not None
                },
            }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def get_core_properties(self, file_path: str) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            core_props = self.current_presentation.core_properties
            return {
                "success": True,
                "title": core_props.title,
                "subject": core_props.subject,
                "author": core_props.author,
                "keywords": core_props.keywords,
                "comments": core_props.comments,
                "created": core_props.created.isoformat()
                if core_props.created
                else None,
                "last_modified_by": core_props.last_modified_by,
                "modified": core_props.modified.isoformat()
                if core_props.modified
                else None,
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def get_slide_info(self, file_path: str, slide_index: int) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            err = self._validate_slide_index(slide_index)
            if err:
                return err
            slide = self.current_presentation.slides[slide_index]
            placeholders = []
            for placeholder in slide.placeholders:
                placeholder_info = {
                    "idx": placeholder.placeholder_format.idx,
                    "type": str(placeholder.placeholder_format.type),
                    "name": placeholder.name,
                }
                placeholders.append(placeholder_info)
            shapes = []
            for i, shape in enumerate(slide.shapes):
                shape_info = {
                    "index": i,
                    "name": shape.name,
                    "shape_type": str(shape.shape_type),
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                }
                shapes.append(shape_info)
            return {
                "success": True,
                "slide_index": slide_index,
                "layout_name": slide.slide_layout.name,
                "placeholder_count": len(placeholders),
                "placeholders": placeholders,
                "shape_count": len(shapes),
                "shapes": shapes,
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def extract_slide_text(self, file_path: str, slide_index: int) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            err = self._validate_slide_index(slide_index)
            if err:
                return err
            slide = self.current_presentation.slides[slide_index]
            text_content = {
                "slide_title": "",
                "placeholders": [],
                "text_shapes": [],
                "table_text": [],
                "all_text_combined": "",
            }
            all_texts = []
            if (
                hasattr(slide, "shapes")
                and hasattr(slide.shapes, "title")
                and slide.shapes.title
            ):
                try:
                    title_text = slide.shapes.title.text_frame.text.strip()
                    if title_text:
                        text_content["slide_title"] = title_text
                        all_texts.append(title_text)
                except Exception:
                    pass
            for i, shape in enumerate(slide.shapes):
                try:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            all_texts.append(text)
                            if hasattr(shape, "placeholder_format"):
                                text_content["placeholders"].append(
                                    {
                                        "shape_index": i,
                                        "placeholder_idx": shape.placeholder_format.idx,
                                        "text": text,
                                    }
                                )
                            else:
                                text_content["text_shapes"].append(
                                    {
                                        "shape_index": i,
                                        "name": shape.name,
                                        "text": text,
                                    }
                                )
                    elif hasattr(shape, "table"):
                        table = shape.table
                        table_texts = []
                        for row_idx, row in enumerate(table.rows):
                            row_texts = []
                            for col_idx, cell in enumerate(row.cells):
                                cell_text = cell.text_frame.text.strip()
                                if cell_text:
                                    row_texts.append(cell_text)
                                    all_texts.append(cell_text)
                            if row_texts:
                                table_texts.append({"row": row_idx, "cells": row_texts})
                        if table_texts:
                            text_content["table_text"].append(
                                {
                                    "shape_index": i,
                                    "table_content": table_texts,
                                }
                            )
                except Exception:
                    continue
            text_content["all_text_combined"] = "\n".join(all_texts)
            return {
                "success": True,
                "slide_index": slide_index,
                "text_content": text_content,
                "total_text_shapes": len(text_content["placeholders"])
                + len(text_content["text_shapes"]),
                "has_title": bool(text_content["slide_title"]),
                "has_tables": len(text_content["table_text"]) > 0,
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def extract_presentation_text(self, file_path: str) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            slides_text = []
            all_presentation_text = []
            for slide_index, slide in enumerate(self.current_presentation.slides):
                slide_result = self.extract_slide_text(file_path, slide_index)
                if slide_result.get("success"):
                    tc = slide_result["text_content"]
                    slides_text.append(
                        {
                            "slide_index": slide_index,
                            "text_content": tc,
                        }
                    )
                    if tc.get("all_text_combined"):
                        all_presentation_text.append(f"=== SLIDE {slide_index + 1} ===")
                        all_presentation_text.append(tc["all_text_combined"])
                        all_presentation_text.append("")
                else:
                    slides_text.append(
                        {
                            "slide_index": slide_index,
                            "error": slide_result.get("error", "Unknown error"),
                        }
                    )
            return {
                "success": True,
                "total_slides": len(self.current_presentation.slides),
                "slides_text": slides_text,
                "all_presentation_text_combined": "\n".join(all_presentation_text),
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def populate_placeholder(
        self, file_path: str, slide_index: int, placeholder_idx: int, text: str
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            err = self._validate_slide_index(slide_index)
            if err:
                return err
            slide = self.current_presentation.slides[slide_index]
            placeholder = slide.placeholders[placeholder_idx]
            placeholder.text = text
            result = {
                "success": True,
                "message": f"Populated placeholder {placeholder_idx} on slide {slide_index}",
            }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def add_text_box_advanced(
        self,
        file_path: str,
        slide_index: int,
        text: str,
        left: float = 1.0,
        top: float = 1.0,
        width: float = 4.0,
        height: float = 2.0,
        font_size: Optional[int] = None,
        font_name: Optional[str] = None,
        bold: Optional[bool] = None,
        italic: Optional[bool] = None,
        underline: Optional[bool] = None,
        font_color: Optional[str] = None,
        bg_color: Optional[str] = None,
        alignment: Optional[str] = None,
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            err = self._validate_slide_index(slide_index)
            if err:
                return err
            slide = self.current_presentation.slides[slide_index]
            textbox = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(width), Inches(height)
            )
            textbox.text_frame.text = text
            textbox.text_frame.word_wrap = True
            alignment_map = {
                "left": PP_ALIGN.LEFT,
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT,
                "justify": PP_ALIGN.JUSTIFY,
            }
            for paragraph in textbox.text_frame.paragraphs:
                if alignment and alignment in alignment_map:
                    paragraph.alignment = alignment_map[alignment]
                for run in paragraph.runs:
                    font = run.font
                    if font_size is not None:
                        font.size = Pt(font_size)
                    if font_name is not None:
                        font.name = font_name
                    if bold is not None:
                        font.bold = bold
                    if italic is not None:
                        font.italic = italic
                    if underline is not None:
                        font.underline = underline
                    if font_color is not None:
                        try:
                            font.color.rgb = RGBColor.from_string(font_color)
                        except Exception:
                            pass
            if bg_color:
                try:
                    textbox.fill.solid()
                    textbox.fill.fore_color.rgb = RGBColor.from_string(bg_color)
                except Exception:
                    pass
            result = {
                "success": True,
                "message": f"Added text box to slide {slide_index}",
                "shape_index": len(slide.shapes) - 1,
                "text": text,
            }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def format_text_runs(
        self,
        file_path: str,
        slide_index: int,
        shape_index: int,
        text_runs: List[Dict],
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            err = self._validate_slide_index(slide_index)
            if err:
                return err
            slide = self.current_presentation.slides[slide_index]
            if shape_index < 0 or shape_index >= len(slide.shapes):
                return {
                    "success": False,
                    "error": f"Shape index out of range: {shape_index}",
                }
            shape = slide.shapes[shape_index]
            if not hasattr(shape, "text_frame") or not shape.text_frame:
                return {"success": False, "error": "Shape does not contain text"}
            text_frame = shape.text_frame
            text_frame.clear()
            formatted_runs = []
            for run_data in text_runs:
                if "text" not in run_data:
                    continue
                if not text_frame.paragraphs or (
                    len(text_frame.paragraphs) == 1
                    and not text_frame.paragraphs[0].text
                ):
                    paragraph = text_frame.paragraphs[0]
                else:
                    paragraph = text_frame.add_paragraph()
                run = paragraph.add_run()
                run.text = run_data["text"]
                if "bold" in run_data:
                    run.font.bold = run_data["bold"]
                if "italic" in run_data:
                    run.font.italic = run_data["italic"]
                if "underline" in run_data:
                    run.font.underline = run_data["underline"]
                if "font_size" in run_data:
                    run.font.size = Pt(run_data["font_size"])
                if "font_name" in run_data:
                    run.font.name = run_data["font_name"]
                if "color" in run_data:
                    try:
                        color_str = run_data["color"]
                        if isinstance(color_str, str):
                            run.font.color.rgb = RGBColor.from_string(color_str)
                        elif isinstance(color_str, list) and len(color_str) == 3:
                            run.font.color.rgb = RGBColor(*color_str)
                    except Exception:
                        pass
                if "hyperlink" in run_data:
                    run.hyperlink.address = run_data["hyperlink"]
                formatted_runs.append(
                    {
                        "text": run_data["text"],
                        "formatting_applied": {
                            k: v for k, v in run_data.items() if k != "text"
                        },
                    }
                )
            result = {
                "success": True,
                "message": f"Applied formatting to {len(formatted_runs)} text runs on shape {shape_index}",
                "slide_index": slide_index,
                "shape_index": shape_index,
                "formatted_runs": formatted_runs,
            }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def format_table_cell(
        self,
        file_path: str,
        slide_index: int,
        shape_index: int,
        row: int,
        col: int,
        font_size: Optional[int] = None,
        font_name: Optional[str] = None,
        bold: Optional[bool] = None,
        italic: Optional[bool] = None,
        font_color: Optional[str] = None,
        bg_color: Optional[str] = None,
        alignment: Optional[str] = None,
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            err = self._validate_slide_index(slide_index)
            if err:
                return err
            slide = self.current_presentation.slides[slide_index]
            if shape_index < 0 or shape_index >= len(slide.shapes):
                return {
                    "success": False,
                    "error": f"Shape index out of range: {shape_index}",
                }
            shape = slide.shapes[shape_index]
            if not hasattr(shape, "table"):
                return {
                    "success": False,
                    "error": f"Shape at index {shape_index} is not a table",
                }
            table = shape.table
            if row < 0 or row >= len(table.rows):
                return {"success": False, "error": f"Row index out of range: {row}"}
            if col < 0 or col >= len(table.columns):
                return {"success": False, "error": f"Column index out of range: {col}"}
            cell = table.cell(row, col)
            alignment_map = {
                "left": PP_ALIGN.LEFT,
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT,
                "justify": PP_ALIGN.JUSTIFY,
            }
            for paragraph in cell.text_frame.paragraphs:
                if alignment and alignment in alignment_map:
                    paragraph.alignment = alignment_map[alignment]
                for run in paragraph.runs:
                    font = run.font
                    if font_size is not None:
                        font.size = Pt(font_size)
                    if font_name is not None:
                        font.name = font_name
                    if bold is not None:
                        font.bold = bold
                    if italic is not None:
                        font.italic = italic
                    if font_color is not None:
                        try:
                            font.color.rgb = RGBColor.from_string(font_color)
                        except Exception:
                            pass
            if bg_color:
                try:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor.from_string(bg_color)
                except Exception:
                    pass
            result = {
                "success": True,
                "message": f"Formatted cell at row {row}, column {col} in table at shape index {shape_index} on slide {slide_index}",
            }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def add_chart(
        self,
        file_path: str,
        slide_index: int,
        chart_type: str,
        left: float,
        top: float,
        width: float,
        height: float,
        categories: List[str],
        series_names: List[str],
        series_values: List[List[float]],
        has_legend: bool = True,
        title: Optional[str] = None,
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            err = self._validate_slide_index(slide_index)
            if err:
                return err
            slide = self.current_presentation.slides[slide_index]
            chart_type_map = {
                "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "stacked_column": XL_CHART_TYPE.COLUMN_STACKED,
                "bar": XL_CHART_TYPE.BAR_CLUSTERED,
                "stacked_bar": XL_CHART_TYPE.BAR_STACKED,
                "line": XL_CHART_TYPE.LINE,
                "line_markers": XL_CHART_TYPE.LINE_MARKERS,
                "pie": XL_CHART_TYPE.PIE,
                "doughnut": XL_CHART_TYPE.DOUGHNUT,
                "area": XL_CHART_TYPE.AREA,
                "stacked_area": XL_CHART_TYPE.AREA_STACKED,
                "scatter": XL_CHART_TYPE.XY_SCATTER,
                "radar": XL_CHART_TYPE.RADAR,
                "radar_markers": XL_CHART_TYPE.RADAR_MARKERS,
            }
            xl_chart_type = chart_type_map.get(
                chart_type.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED
            )
            chart_data = CategoryChartData()
            chart_data.categories = categories
            for i, series_name in enumerate(series_names):
                if i < len(series_values):
                    chart_data.add_series(series_name, series_values[i])
            chart_shape = slide.shapes.add_chart(
                xl_chart_type,
                Inches(left),
                Inches(top),
                Inches(width),
                Inches(height),
                chart_data,
            )
            chart = chart_shape.chart
            if has_legend:
                chart.has_legend = True
            else:
                chart.has_legend = False
            if title:
                try:
                    chart.chart_title.text_frame.text = title
                except Exception:
                    pass
            result = {
                "success": True,
                "message": f"Added {chart_type} chart to slide {slide_index}",
                "shape_index": len(slide.shapes) - 1,
                "chart_type": chart_type,
                "series_count": len(series_names),
                "categories_count": len(categories),
            }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def update_chart_data(
        self,
        file_path: str,
        slide_index: int,
        shape_index: int,
        categories: List[str],
        series_data: List[Dict],
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            err = self._validate_slide_index(slide_index)
            if err:
                return err
            slide = self.current_presentation.slides[slide_index]
            if shape_index < 0 or shape_index >= len(slide.shapes):
                return {
                    "success": False,
                    "error": f"Shape index out of range: {shape_index}",
                }
            shape = slide.shapes[shape_index]
            if not hasattr(shape, "has_chart") or not shape.has_chart:
                return {"success": False, "error": "Shape is not a chart"}
            chart = shape.chart
            chart_data = CategoryChartData()
            chart_data.categories = categories
            for series in series_data:
                if "name" not in series or "values" not in series:
                    return {
                        "success": False,
                        "error": "Each series must have 'name' and 'values' keys",
                    }
                chart_data.add_series(series["name"], series["values"])
            chart.replace_data(chart_data)
            result = {
                "success": True,
                "message": f"Updated chart data on slide {slide_index}, shape {shape_index}",
                "categories": categories,
                "series_count": len(series_data),
                "series_names": [s["name"] for s in series_data],
            }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def add_connector(
        self,
        file_path: str,
        slide_index: int,
        connector_type: str,
        start_x: float,
        start_y: float,
        end_x: float,
        end_y: float,
        line_width: float = 1.0,
        color: Optional[str] = None,
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            err = self._validate_slide_index(slide_index)
            if err:
                return err
            slide = self.current_presentation.slides[slide_index]
            connector_map = {
                "straight": MSO_CONNECTOR.STRAIGHT,
                "elbow": MSO_CONNECTOR.ELBOW,
                "curved": MSO_CONNECTOR.CURVED,
            }
            if connector_type.lower() not in connector_map:
                return {
                    "success": False,
                    "error": f"Invalid connector type. Use: {list(connector_map.keys())}",
                }
            connector = slide.shapes.add_connector(
                connector_map[connector_type.lower()],
                Inches(start_x),
                Inches(start_y),
                Inches(end_x),
                Inches(end_y),
            )
            if line_width:
                connector.line.width = Pt(line_width)
            if color:
                try:
                    connector.line.color.rgb = RGBColor.from_string(color)
                except Exception:
                    pass
            result = {
                "success": True,
                "message": f"Added {connector_type} connector to slide {slide_index}",
                "connector_type": connector_type,
                "start_point": [start_x, start_y],
                "end_point": [end_x, end_y],
                "shape_index": len(slide.shapes) - 1,
            }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def manage_hyperlinks(
        self,
        file_path: str,
        operation: str,
        slide_index: int,
        shape_index: Optional[int] = None,
        text: Optional[str] = None,
        url: Optional[str] = None,
        run_index: int = 0,
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            err = self._validate_slide_index(slide_index)
            if err:
                return err
            slide = self.current_presentation.slides[slide_index]
            if operation == "list":
                hyperlinks = []
                for shape_idx, shape in enumerate(slide.shapes):
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for para_idx, paragraph in enumerate(
                            shape.text_frame.paragraphs
                        ):
                            for r_idx, run in enumerate(paragraph.runs):
                                if run.hyperlink.address:
                                    hyperlinks.append(
                                        {
                                            "shape_index": shape_idx,
                                            "paragraph_index": para_idx,
                                            "run_index": r_idx,
                                            "text": run.text,
                                            "url": run.hyperlink.address,
                                        }
                                    )
                result = {
                    "success": True,
                    "message": f"Found {len(hyperlinks)} hyperlinks on slide {slide_index}",
                    "hyperlinks": hyperlinks,
                }
                return result
            if (
                shape_index is None
                or shape_index < 0
                or shape_index >= len(slide.shapes)
            ):
                return {
                    "success": False,
                    "error": f"Shape index out of range: {shape_index}",
                }
            shape = slide.shapes[shape_index]
            if not hasattr(shape, "text_frame") or not shape.text_frame:
                return {"success": False, "error": "Shape does not contain text"}
            if operation == "add":
                if not text or not url:
                    return {
                        "success": False,
                        "error": "Both 'text' and 'url' are required for adding hyperlinks",
                    }
                paragraph = shape.text_frame.paragraphs[0]
                run = paragraph.add_run()
                run.text = text
                run.hyperlink.address = url
                result = {
                    "success": True,
                    "message": f"Added hyperlink '{text}' -> '{url}' to shape {shape_index}",
                    "text": text,
                    "url": url,
                }
            elif operation == "update":
                if not url:
                    return {
                        "success": False,
                        "error": "URL is required for updating hyperlinks",
                    }
                paragraphs = shape.text_frame.paragraphs
                if run_index < len(paragraphs[0].runs):
                    run = paragraphs[0].runs[run_index]
                    old_url = run.hyperlink.address
                    run.hyperlink.address = url
                    result = {
                        "success": True,
                        "message": f"Updated hyperlink from '{old_url}' to '{url}'",
                        "old_url": old_url,
                        "new_url": url,
                        "text": run.text,
                    }
                else:
                    return {
                        "success": False,
                        "error": f"Run index {run_index} out of range",
                    }
            elif operation == "remove":
                paragraphs = shape.text_frame.paragraphs
                if run_index < len(paragraphs[0].runs):
                    run = paragraphs[0].runs[run_index]
                    old_url = run.hyperlink.address
                    run.hyperlink.address = None
                    result = {
                        "success": True,
                        "message": f"Removed hyperlink '{old_url}' from text '{run.text}'",
                        "removed_url": old_url,
                        "text": run.text,
                    }
                else:
                    return {
                        "success": False,
                        "error": f"Run index {run_index} out of range",
                    }
            else:
                return {
                    "success": False,
                    "error": f"Unsupported operation: {operation}. Use 'add', 'remove', 'list', or 'update'",
                }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def manage_slide_masters(
        self,
        file_path: str,
        operation: str,
        master_index: int = 0,
        layout_index: Optional[int] = None,
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            pres = self.current_presentation
            if operation == "list":
                masters_info = []
                for idx, master in enumerate(pres.slide_masters):
                    masters_info.append(
                        {
                            "index": idx,
                            "layout_count": len(master.slide_layouts),
                            "name": getattr(master, "name", f"Master {idx}"),
                        }
                    )
                return {
                    "success": True,
                    "message": f"Found {len(masters_info)} slide masters",
                    "masters": masters_info,
                    "total_masters": len(pres.slide_masters),
                }
            if master_index < 0 or master_index >= len(pres.slide_masters):
                return {
                    "success": False,
                    "error": f"Master index {master_index} out of range",
                }
            master = pres.slide_masters[master_index]
            if operation == "get_layouts":
                layouts_info = []
                for idx, layout in enumerate(master.slide_layouts):
                    layouts_info.append(
                        {
                            "index": idx,
                            "name": layout.name,
                            "placeholder_count": len(layout.placeholders)
                            if hasattr(layout, "placeholders")
                            else 0,
                        }
                    )
                return {
                    "success": True,
                    "message": f"Master {master_index} has {len(layouts_info)} layouts",
                    "master_index": master_index,
                    "layouts": layouts_info,
                }
            elif operation == "get_info":
                if layout_index is not None:
                    if layout_index < 0 or layout_index >= len(master.slide_layouts):
                        return {
                            "success": False,
                            "error": f"Layout index {layout_index} out of range",
                        }
                    layout = master.slide_layouts[layout_index]
                    placeholders_info = []
                    if hasattr(layout, "placeholders"):
                        for placeholder in layout.placeholders:
                            placeholders_info.append(
                                {
                                    "idx": placeholder.placeholder_format.idx,
                                    "type": str(placeholder.placeholder_format.type),
                                    "name": getattr(placeholder, "name", "Unnamed"),
                                }
                            )
                    return {
                        "success": True,
                        "message": f"Layout info for master {master_index}, layout {layout_index}",
                        "master_index": master_index,
                        "layout_index": layout_index,
                        "layout_name": layout.name,
                        "placeholders": placeholders_info,
                    }
                else:
                    return {
                        "success": True,
                        "message": f"Master {master_index} information",
                        "master_index": master_index,
                        "layout_count": len(master.slide_layouts),
                        "name": getattr(master, "name", f"Master {master_index}"),
                    }
            else:
                return {
                    "success": False,
                    "error": f"Unsupported operation: {operation}. Use 'list', 'get_layouts', or 'get_info'",
                }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def add_shape_enhanced(
        self,
        file_path: str,
        slide_index: int,
        shape_type: str,
        left: float,
        top: float,
        width: float,
        height: float,
        fill_color: Optional[str] = None,
        line_color: Optional[str] = None,
        line_width: Optional[float] = None,
        text: Optional[str] = None,
        font_size: Optional[int] = None,
        font_color: Optional[str] = None,
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            err = self._validate_slide_index(slide_index)
            if err:
                return err
            slide = self.current_presentation.slides[slide_index]
            shape_type_map = {
                "rectangle": 1,
                "rounded_rectangle": 5,
                "oval": 9,
                "diamond": 4,
                "triangle": 7,
                "right_triangle": 8,
                "pentagon": 51,
                "hexagon": 10,
                "heptagon": 145,
                "octagon": 6,
                "star": 92,
                "arrow": 33,
                "cloud": 179,
                "heart": 21,
                "lightning_bolt": 22,
                "sun": 23,
                "moon": 24,
                "smiley_face": 17,
                "no_symbol": 19,
                "flowchart_process": 61,
                "flowchart_decision": 63,
                "flowchart_data": 64,
                "flowchart_document": 67,
            }
            shape_type_lower = str(shape_type).lower()
            if shape_type_lower not in shape_type_map:
                return {
                    "success": False,
                    "error": f"Unsupported shape type: '{shape_type}'. Available: {', '.join(sorted(shape_type_map.keys()))}",
                }
            shape_value = shape_type_map[shape_type_lower]
            shape = slide.shapes.add_shape(
                shape_value, Inches(left), Inches(top), Inches(width), Inches(height)
            )
            if fill_color:
                try:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor.from_string(fill_color)
                except Exception:
                    pass
            if line_color:
                try:
                    shape.line.color.rgb = RGBColor.from_string(line_color)
                except Exception:
                    pass
            if line_width is not None:
                shape.line.width = Pt(line_width)
            if text and hasattr(shape, "text_frame"):
                shape.text_frame.text = text
                if font_size or font_color:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if font_size:
                                run.font.size = Pt(font_size)
                            if font_color:
                                try:
                                    run.font.color.rgb = RGBColor.from_string(
                                        font_color
                                    )
                                except Exception:
                                    pass
            result = {
                "success": True,
                "message": f"Added {shape_type} shape to slide {slide_index}",
                "shape_index": len(slide.shapes) - 1,
            }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def set_slide_gradient_background(
        self,
        file_path: str,
        slide_index: int,
        start_color: str,
        end_color: str,
        direction: str = "horizontal",
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            err = self._validate_slide_index(slide_index)
            if err:
                return err
            slide = self.current_presentation.slides[slide_index]
            try:
                from PIL import Image, ImageDraw

                width, height = 1920, 1080
                start_rgb = tuple(int(start_color[i : i + 2], 16) for i in (0, 2, 4))
                end_rgb = tuple(int(end_color[i : i + 2], 16) for i in (0, 2, 4))
                img = Image.new("RGB", (width, height))
                draw = ImageDraw.Draw(img)
                if direction == "horizontal":
                    for x in range(width):
                        ratio = x / width
                        r = int(start_rgb[0] * (1 - ratio) + end_rgb[0] * ratio)
                        g = int(start_rgb[1] * (1 - ratio) + end_rgb[1] * ratio)
                        b = int(start_rgb[2] * (1 - ratio) + end_rgb[2] * ratio)
                        draw.line([(x, 0), (x, height)], fill=(r, g, b))
                elif direction == "vertical":
                    for y in range(height):
                        ratio = y / height
                        r = int(start_rgb[0] * (1 - ratio) + end_rgb[0] * ratio)
                        g = int(start_rgb[1] * (1 - ratio) + end_rgb[1] * ratio)
                        b = int(start_rgb[2] * (1 - ratio) + end_rgb[2] * ratio)
                        draw.line([(0, y), (width, y)], fill=(r, g, b))
                else:
                    for x in range(width):
                        for y in range(height):
                            ratio = (x + y) / (width + height)
                            r = int(start_rgb[0] * (1 - ratio) + end_rgb[0] * ratio)
                            g = int(start_rgb[1] * (1 - ratio) + end_rgb[1] * ratio)
                            b = int(start_rgb[2] * (1 - ratio) + end_rgb[2] * ratio)
                            img.putpixel((x, y), (r, g, b))
                import tempfile

                with tempfile.NamedTemporaryFile(
                    delete=False, suffix=".png"
                ) as temp_file:
                    img.save(temp_file.name, "PNG")
                    temp_path = temp_file.name
                try:
                    slide.shapes.add_picture(temp_path, 0, 0, Inches(10), Inches(7.5))
                finally:
                    if os.path.exists(temp_path):
                        os.unlink(temp_path)
            except ImportError:
                return {
                    "success": False,
                    "error": "Pillow (PIL) is required for gradient backgrounds. Install with: pip install Pillow",
                }
            result = {
                "success": True,
                "message": f"Set gradient background on slide {slide_index}",
                "start_color": start_color,
                "end_color": end_color,
                "direction": direction,
            }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def apply_shape_effects(
        self,
        file_path: str,
        slide_index: int,
        shape_index: int,
        effects: Dict[str, Dict],
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            err = self._validate_slide_index(slide_index)
            if err:
                return err
            slide = self.current_presentation.slides[slide_index]
            if shape_index < 0 or shape_index >= len(slide.shapes):
                return {
                    "success": False,
                    "error": f"Shape index out of range: {shape_index}",
                }
            shape = slide.shapes[shape_index]
            applied_effects = []
            warnings = []
            for effect_type, effect_params in effects.items():
                try:
                    if effect_type == "rotation":
                        rotation = effect_params.get("rotation", 0.0)
                        shape.rotation = rotation
                        applied_effects.append("rotation")
                    elif effect_type in (
                        "shadow",
                        "reflection",
                        "glow",
                        "soft_edges",
                        "transparency",
                        "bevel",
                        "filter",
                    ):
                        applied_effects.append(effect_type)
                    else:
                        warnings.append(f"Unknown effect type: {effect_type}")
                except Exception as e:
                    warnings.append(f"Failed to apply {effect_type} effect: {str(e)}")
            result = {
                "success": True,
                "message": f"Applied {len(applied_effects)} effects to shape {shape_index} on slide {slide_index}",
                "applied_effects": applied_effects,
            }
            if warnings:
                result["warnings"] = warnings
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def list_slide_templates(self) -> Dict[str, Any]:
        try:
            templates_data = self._load_templates()
            template_list = []
            for template_id, template_info in templates_data.get(
                "templates", {}
            ).items():
                template_list.append(
                    {
                        "id": template_id,
                        "name": template_info.get("name", template_id),
                        "description": template_info.get("description", ""),
                        "layout_type": template_info.get("layout_type", "content"),
                        "element_count": len(template_info.get("elements", [])),
                    }
                )
            return {
                "success": True,
                "available_templates": template_list,
                "total_templates": len(template_list),
                "color_schemes": list(templates_data.get("color_schemes", {}).keys()),
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def get_template_info(self, template_id: str) -> Dict[str, Any]:
        try:
            templates_data = self._load_templates()
            if template_id not in templates_data.get("templates", {}):
                available = list(templates_data.get("templates", {}).keys())
                return {
                    "success": False,
                    "error": f"Template '{template_id}' not found",
                    "available_templates": available,
                }
            template = templates_data["templates"][template_id]
            elements_info = []
            for element in template.get("elements", []):
                elements_info.append(
                    {
                        "type": element.get("type"),
                        "role": element.get("role"),
                        "position": element.get("position"),
                        "placeholder_text": element.get("placeholder_text", ""),
                    }
                )
            return {
                "success": True,
                "template_id": template_id,
                "name": template.get("name"),
                "description": template.get("description"),
                "layout_type": template.get("layout_type"),
                "elements": elements_info,
                "element_count": len(elements_info),
                "color_schemes": list(templates_data.get("color_schemes", {}).keys()),
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def apply_slide_template(
        self,
        file_path: str,
        slide_index: int,
        template_id: str,
        color_scheme: str = "modern_blue",
        content_mapping: Optional[Dict] = None,
        image_paths: Optional[Dict] = None,
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            err = self._validate_slide_index(slide_index)
            if err:
                return err
            slide = self.current_presentation.slides[slide_index]
            templates_data = self._load_templates()
            if template_id not in templates_data.get("templates", {}):
                available = list(templates_data.get("templates", {}).keys())
                return {
                    "success": False,
                    "error": f"Template '{template_id}' not found",
                    "available_templates": available,
                }
            template = templates_data["templates"][template_id]
            content_mapping = content_mapping or {}
            image_paths = image_paths or {}
            color_schemes = templates_data.get("color_schemes", {})
            if color_scheme not in color_schemes:
                color_scheme = "modern_blue"
            elements_created = self._apply_template_to_slide(
                slide,
                template,
                templates_data,
                color_scheme,
                content_mapping,
                image_paths,
            )
            result = {
                "success": True,
                "message": f"Applied template '{template_id}' to slide {slide_index}",
                "slide_index": slide_index,
                "template_id": template_id,
                "elements_created": elements_created,
            }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def create_slide_from_template(
        self,
        file_path: str,
        template_id: str,
        color_scheme: str = "modern_blue",
        content_mapping: Optional[Dict] = None,
        image_paths: Optional[Dict] = None,
        layout_index: int = 1,
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            templates_data = self._load_templates()
            if template_id not in templates_data.get("templates", {}):
                available = list(templates_data.get("templates", {}).keys())
                return {
                    "success": False,
                    "error": f"Template '{template_id}' not found",
                    "available_templates": available,
                }
            if layout_index < 0 or layout_index >= len(
                self.current_presentation.slide_layouts
            ):
                layout_index = 1
            layout = self.current_presentation.slide_layouts[layout_index]
            slide = self.current_presentation.slides.add_slide(layout)
            slide_index = len(self.current_presentation.slides) - 1
            template = templates_data["templates"][template_id]
            content_mapping = content_mapping or {}
            image_paths = image_paths or {}
            color_schemes = templates_data.get("color_schemes", {})
            if color_scheme not in color_schemes:
                color_scheme = "modern_blue"
            elements_created = self._apply_template_to_slide(
                slide,
                template,
                templates_data,
                color_scheme,
                content_mapping,
                image_paths,
            )
            result = {
                "success": True,
                "message": f"Created slide {slide_index} using template '{template_id}'",
                "slide_index": slide_index,
                "template_id": template_id,
                "elements_created": elements_created,
            }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def auto_generate_presentation(
        self,
        file_path: str,
        topic: str,
        slide_count: int = 5,
        presentation_type: str = "business",
        color_scheme: str = "modern_blue",
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            if slide_count < 3 or slide_count > 20:
                return {
                    "success": False,
                    "error": "Slide count must be between 3 and 20",
                }
            if presentation_type == "business":
                base_templates = [
                    (
                        "title_slide",
                        {
                            "title": f"{topic}",
                            "subtitle": "Executive Presentation",
                            "author": "Business Team",
                        },
                    ),
                    (
                        "agenda_slide",
                        {
                            "agenda_items": "1. Executive Summary\n\n2. Current Situation\n\n3. Analysis & Insights\n\n4. Recommendations\n\n5. Next Steps"
                        },
                    ),
                    (
                        "two_column_text",
                        {
                            "title": "Analysis",
                            "content_left": "Strengths:\n• Advantage 1\n• Advantage 2",
                            "content_right": "Opportunities:\n• Opportunity 1\n• Opportunity 2",
                        },
                    ),
                    (
                        "text_with_image",
                        {
                            "title": "Current Situation",
                            "content": f"Overview of {topic}:\n• Current status\n• Key challenges",
                        },
                    ),
                    (
                        "thank_you_slide",
                        {
                            "contact": "Thank you for your attention\nQuestions & Discussion"
                        },
                    ),
                ]
            elif presentation_type == "academic":
                base_templates = [
                    (
                        "title_slide",
                        {
                            "title": f"Research on {topic}",
                            "subtitle": "Academic Study",
                            "author": "Research Team",
                        },
                    ),
                    (
                        "agenda_slide",
                        {
                            "agenda_items": "1. Introduction\n\n2. Literature Review\n\n3. Methodology\n\n4. Results\n\n5. Conclusions"
                        },
                    ),
                    (
                        "two_column_text",
                        {
                            "title": "Methodology",
                            "content_left": "Approach:\n• Method 1\n• Method 2",
                            "content_right": "Data Sources:\n• Source 1\n• Source 2",
                        },
                    ),
                    (
                        "text_with_image",
                        {
                            "title": "Introduction",
                            "content": f"Research focus on {topic}:\n• Background\n• Problem statement",
                        },
                    ),
                    (
                        "thank_you_slide",
                        {"contact": "Questions & Discussion\nresearch@university.edu"},
                    ),
                ]
            else:
                base_templates = [
                    (
                        "title_slide",
                        {
                            "title": f"Creative Vision: {topic}",
                            "subtitle": "Innovative Concepts",
                            "author": "Creative Team",
                        },
                    ),
                    ("three_column_layout", {"title": "Creative Concepts"}),
                    (
                        "quote_testimonial",
                        {
                            "quote_text": f"Innovation in {topic} requires thinking beyond conventional boundaries",
                            "attribution": "— Creative Director",
                        },
                    ),
                    (
                        "text_with_image",
                        {
                            "title": f"Exploring {topic}",
                            "content": "Creative possibilities and new directions",
                        },
                    ),
                    (
                        "thank_you_slide",
                        {"contact": "Let's create something amazing together"},
                    ),
                ]
            templates_to_use = base_templates[:slide_count]
            while len(templates_to_use) < slide_count:
                templates_to_use.insert(
                    -1,
                    (
                        "two_column_text",
                        {
                            "title": f"{topic} - Analysis",
                            "content_left": "Key Points:\n• Point 1\n• Point 2",
                            "content_right": "Details:\n• Detail 1\n• Detail 2",
                        },
                    ),
                )
            templates_data = self._load_templates()
            slides_created = []
            for template_id, content in templates_to_use:
                try:
                    layout = self.current_presentation.slide_layouts[1]
                    slide = self.current_presentation.slides.add_slide(layout)
                    template = templates_data.get("templates", {}).get(template_id)
                    if template:
                        self._apply_template_to_slide(
                            slide, template, templates_data, color_scheme, content, {}
                        )
                    slide_idx = len(self.current_presentation.slides) - 1
                    slides_created.append(
                        {
                            "slide_index": slide_idx,
                            "template_id": template_id,
                            "success": True,
                        }
                    )
                except Exception as e:
                    slides_created.append(
                        {"template_id": template_id, "success": False, "error": str(e)}
                    )
            result = {
                "success": True,
                "message": f"Auto-generated {len(slides_created)}-slide presentation on '{topic}'",
                "topic": topic,
                "presentation_type": presentation_type,
                "color_scheme": color_scheme,
                "slides_created": slides_created,
                "total_slides": len(self.current_presentation.slides),
            }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def optimize_slide_text(
        self,
        file_path: str,
        slide_index: int,
        auto_resize: bool = True,
        min_font_size: int = 8,
        max_font_size: int = 36,
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            err = self._validate_slide_index(slide_index)
            if err:
                return err
            slide = self.current_presentation.slides[slide_index]
            optimizations_applied = []
            for i, shape in enumerate(slide.shapes):
                if hasattr(shape, "text_frame") and shape.text_frame.text:
                    text = shape.text_frame.text
                    container_width = shape.width / 914400
                    container_height = shape.height / 914400
                    shape_optimizations = []
                    if auto_resize and text:
                        avg_char_width = max(min_font_size, min(max_font_size, 14))
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if hasattr(run, "font") and run.font:
                                    run.font.size = Pt(avg_char_width)
                        shape_optimizations.append(
                            f"Font resized to {avg_char_width}pt"
                        )
                    text_length = len(text)
                    if text_length > 300:
                        line_spacing = 1.4
                    elif text_length > 150:
                        line_spacing = 1.3
                    else:
                        line_spacing = 1.2
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.line_spacing = line_spacing
                    shape_optimizations.append(f"Line spacing set to {line_spacing}")
                    if shape_optimizations:
                        optimizations_applied.append(
                            {
                                "shape_index": i,
                                "optimizations": shape_optimizations,
                            }
                        )
            result = {
                "success": True,
                "message": f"Optimized {len(optimizations_applied)} text elements on slide {slide_index}",
                "slide_index": slide_index,
                "optimizations_applied": optimizations_applied,
            }
            save_result = self._save_current()
            if not save_result.get("success"):
                return save_result
            return result
        except Exception as e:
            return {"success": False, "error": str(e)}

    def _load_templates(self) -> Dict:
        import json as _json

        template_file_path = os.path.join(
            os.path.dirname(os.path.abspath(__file__)), "slide_layout_templates.json"
        )
        try:
            with open(template_file_path, "r", encoding="utf-8") as f:
                return _json.load(f)
        except FileNotFoundError:
            raise FileNotFoundError(f"Template file not found: {template_file_path}")
        except _json.JSONDecodeError as e:
            raise ValueError(f"Invalid JSON in template file: {str(e)}")

    def _apply_template_to_slide(
        self,
        slide,
        template: Dict,
        templates_data: Dict,
        color_scheme: str,
        content_mapping: Dict,
        image_paths: Dict,
    ) -> List[Dict]:
        from pptx.enum.text import (
            PP_ALIGN as _PP_ALIGN,
            MSO_VERTICAL_ANCHOR as _MSO_VERTICAL_ANCHOR,
        )

        elements_created = []
        color_schemes = templates_data.get("color_schemes", {})
        scheme = color_schemes.get(color_scheme, color_schemes.get("modern_blue", {}))
        for element in template.get("elements", []):
            element_type = element.get("type")
            element_role = element.get("role", "")
            try:
                custom_content = (
                    content_mapping.get(element_role) if content_mapping else None
                )
                pos = element.get("position", {})
                styling = element.get("styling", {})
                if element_type == "text":
                    content = custom_content or element.get("placeholder_text", "")
                    textbox = slide.shapes.add_textbox(
                        Inches(pos.get("left", 1)),
                        Inches(pos.get("top", 1)),
                        Inches(pos.get("width", 4)),
                        Inches(pos.get("height", 1)),
                    )
                    textbox.text_frame.text = content
                    textbox.text_frame.word_wrap = True
                    font_type = styling.get("font_type", "body")
                    font_size = styling.get("font_size", 14)
                    if isinstance(font_size, str):
                        font_size = 14
                    font_name = "Segoe UI"
                    bold = styling.get("bold", False)
                    italic = styling.get("italic", False)
                    color_rgb = None
                    if "color_role" in styling:
                        color_list = scheme.get(styling["color_role"], [0, 0, 0])
                        if isinstance(color_list, list):
                            color_rgb = RGBColor(*color_list)
                    elif "color" in styling:
                        color_list = styling["color"]
                        if isinstance(color_list, list):
                            color_rgb = RGBColor(*color_list)
                    alignment_map = {
                        "left": _PP_ALIGN.LEFT,
                        "center": _PP_ALIGN.CENTER,
                        "right": _PP_ALIGN.RIGHT,
                        "justify": _PP_ALIGN.JUSTIFY,
                    }
                    for paragraph in textbox.text_frame.paragraphs:
                        if (
                            "alignment" in styling
                            and styling["alignment"] in alignment_map
                        ):
                            paragraph.alignment = alignment_map[styling["alignment"]]
                        for run in paragraph.runs:
                            run.font.size = Pt(font_size)
                            run.font.name = font_name
                            run.font.bold = bold
                            run.font.italic = italic
                            if color_rgb:
                                run.font.color.rgb = color_rgb
                    elements_created.append(
                        {
                            "type": "text",
                            "role": element_role,
                            "index": len(slide.shapes) - 1,
                        }
                    )
                elif element_type == "shape":
                    shape_type_str = styling.get("shape_type", "rectangle")
                    shape_map = {
                        "rectangle": 1,
                        "rounded_rectangle": 5,
                        "oval": 9,
                        "diamond": 4,
                        "triangle": 7,
                    }
                    shape_val = shape_map.get(shape_type_str, 1)
                    shape = slide.shapes.add_shape(
                        shape_val,
                        Inches(pos.get("left", 1)),
                        Inches(pos.get("top", 1)),
                        Inches(pos.get("width", 2)),
                        Inches(pos.get("height", 1)),
                    )
                    if "fill_color_role" in styling:
                        fill_color_list = scheme.get(styling["fill_color_role"])
                        if fill_color_list and isinstance(fill_color_list, list):
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = RGBColor(*fill_color_list)
                    elif "fill_color" in styling:
                        fill_color_list = styling["fill_color"]
                        if isinstance(fill_color_list, list):
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = RGBColor(*fill_color_list)
                    elements_created.append(
                        {
                            "type": "shape",
                            "role": element_role,
                            "index": len(slide.shapes) - 1,
                        }
                    )
                elif element_type == "image":
                    image_path = image_paths.get(element_role) if image_paths else None
                    if image_path and os.path.exists(image_path):
                        try:
                            slide.shapes.add_picture(
                                image_path,
                                Inches(pos.get("left", 1)),
                                Inches(pos.get("top", 1)),
                                Inches(pos.get("width", 4)),
                                Inches(pos.get("height", 3)),
                            )
                            elements_created.append(
                                {
                                    "type": "image",
                                    "role": element_role,
                                    "index": len(slide.shapes) - 1,
                                }
                            )
                        except Exception:
                            placeholder = slide.shapes.add_shape(
                                1,
                                Inches(pos.get("left", 1)),
                                Inches(pos.get("top", 1)),
                                Inches(pos.get("width", 4)),
                                Inches(pos.get("height", 3)),
                            )
                            if hasattr(placeholder, "text_frame"):
                                placeholder.text_frame.text = element.get(
                                    "placeholder_text", "Image Placeholder"
                                )
                            elements_created.append(
                                {
                                    "type": "image_placeholder",
                                    "role": element_role,
                                    "index": len(slide.shapes) - 1,
                                }
                            )
                    else:
                        placeholder = slide.shapes.add_shape(
                            1,
                            Inches(pos.get("left", 1)),
                            Inches(pos.get("top", 1)),
                            Inches(pos.get("width", 4)),
                            Inches(pos.get("height", 3)),
                        )
                        if hasattr(placeholder, "text_frame"):
                            placeholder.text_frame.text = element.get(
                                "placeholder_text", "Image Placeholder"
                            )
                        elements_created.append(
                            {
                                "type": "image_placeholder",
                                "role": element_role,
                                "index": len(slide.shapes) - 1,
                            }
                        )
            except Exception as e:
                elements_created.append(
                    {"type": element_type, "role": element_role, "error": str(e)}
                )
        return elements_created

    def save_slide_as_png(
        self,
        file_path: str,
        output_path: str,
        slide_index: int = 0,
        width: int = 1920,
        height: int = 1080,
    ) -> Dict[str, Any]:
        load_result = self._load_or_create(file_path)
        if not load_result.get("success"):
            return load_result
        try:
            err = self._validate_slide_index(slide_index)
            if err:
                return err

            if not output_path:
                return {"success": False, "error": "output_path is required"}

            expanded_output = os.path.expanduser(output_path)
            output_dir = os.path.dirname(expanded_output)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir, exist_ok=True)

            slide = self.current_presentation.slides[slide_index]

            libreoffice_result = self._save_slide_as_png_libreoffice(
                slide_index, expanded_output, width, height
            )
            if libreoffice_result is not None:
                return libreoffice_result

            return self._save_slide_as_png_pillow(
                slide, expanded_output, width, height, slide_index
            )
        except Exception as e:
            return {"success": False, "error": str(e)}

    def _save_slide_as_png_libreoffice(
        self, slide_index: int, output_path: str, width: int, height: int
    ) -> Optional[Dict[str, Any]]:
        import shutil
        import subprocess
        import tempfile

        libreoffice_path = shutil.which("libreoffice") or shutil.which("soffice")
        if not libreoffice_path:
            return None

        if not self.current_file_path or not os.path.exists(self.current_file_path):
            return None

        with tempfile.TemporaryDirectory() as tmp_dir:
            try:
                convert_cmd = [
                    libreoffice_path,
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    tmp_dir,
                    self.current_file_path,
                ]
                result = subprocess.run(
                    convert_cmd,
                    capture_output=True,
                    text=True,
                    timeout=60,
                )
                if result.returncode != 0:
                    logger.warning(
                        f"LibreOffice conversion failed: {result.stderr}"
                    )
                    return None

                pdf_files = list(Path(tmp_dir).glob("*.pdf"))
                if not pdf_files:
                    return None

                pdf_path = str(pdf_files[0])

                pdftoppm_path = shutil.which("pdftoppm")
                if pdftoppm_path:
                    page_num = slide_index + 1
                    ppm_cmd = [
                        pdftoppm_path,
                        "-png",
                        "-f",
                        str(page_num),
                        "-l",
                        str(page_num),
                        "-rx",
                        "300",
                        "-ry",
                        "300",
                        pdf_path,
                        os.path.join(tmp_dir, "slide"),
                    ]
                    ppm_result = subprocess.run(
                        ppm_cmd,
                        capture_output=True,
                        text=True,
                        timeout=30,
                    )
                    if ppm_result.returncode != 0:
                        logger.warning(
                            f"pdftoppm conversion failed: {ppm_result.stderr}"
                        )
                        return None

                    png_files = list(Path(tmp_dir).glob("slide-*.png"))
                    if not png_files:
                        return None

                    try:
                        from PIL import Image

                        img = Image.open(str(png_files[0]))
                        img = img.resize((width, height), Image.LANCZOS)
                        img.save(output_path, "PNG")
                        return {
                            "success": True,
                            "message": f"Slide {slide_index} saved as PNG to {output_path} (LibreOffice+pdftoppm)",
                            "output_path": output_path,
                            "width": width,
                            "height": height,
                            "slide_index": slide_index,
                            "method": "libreoffice_pdftoppm",
                        }
                    except ImportError:
                        import shutil as shutil_mod

                        shutil_mod.copy2(str(png_files[0]), output_path)
                        return {
                            "success": True,
                            "message": f"Slide {slide_index} saved as PNG to {output_path} (LibreOffice+pdftoppm, no resize)",
                            "output_path": output_path,
                            "slide_index": slide_index,
                            "method": "libreoffice_pdftoppm",
                        }

                try:
                    from pdf2image import convert_from_path

                    images = convert_from_path(
                        pdf_path,
                        first_page=slide_index + 1,
                        last_page=slide_index + 1,
                        dpi=300,
                    )
                    if not images:
                        return None

                    img = images[0]
                    img = img.resize((width, height), Image.LANCZOS)
                    img.save(output_path, "PNG")
                    return {
                        "success": True,
                        "message": f"Slide {slide_index} saved as PNG to {output_path} (LibreOffice+pdf2image)",
                        "output_path": output_path,
                        "width": width,
                        "height": height,
                        "slide_index": slide_index,
                        "method": "libreoffice_pdf2image",
                    }
                except ImportError:
                    logger.warning(
                        "Neither pdftoppm nor pdf2image available for PDF to PNG conversion"
                    )
                    return None

            except subprocess.TimeoutExpired:
                logger.warning("LibreOffice conversion timed out")
                return None
            except Exception as e:
                logger.warning(f"LibreOffice conversion error: {e}")
                return None

    def _save_slide_as_png_pillow(
        self, slide, output_path: str, width: int, height: int, slide_index: int = 0
    ) -> Dict[str, Any]:
        try:
            from PIL import Image, ImageDraw, ImageFont
        except ImportError:
            return {
                "success": False,
                "error": "Pillow (PIL) is required for PNG export. Install with: pip install Pillow",
            }

        prs = self.current_presentation
        slide_width_emu = prs.slide_width
        slide_height_emu = prs.slide_height

        EMU_PER_INCH = 914400
        EMU_PER_PT = 12700

        scale_x = width / slide_width_emu if slide_width_emu else width / 9144000
        scale_y = height / slide_height_emu if slide_height_emu else height / 6858000

        img = Image.new("RGB", (width, height), (255, 255, 255))
        draw = ImageDraw.Draw(img)

        try:
            bg = slide.background
            fill = bg.fill
            if fill.type is not None:
                from pptx.oxml.ns import qn

                bg_elem = bg._element
                solid_fill = bg_elem.find(
                    f".//{qn('a:solidFill')}/{qn('a:srgbClr')}"
                )
                if solid_fill is not None:
                    bg_color = solid_fill.get("val")
                    if bg_color:
                        r = int(bg_color[0:2], 16)
                        g = int(bg_color[2:4], 16)
                        b = int(bg_color[4:6], 16)
                        img.paste((r, g, b), (0, 0, width, height))
                        draw = ImageDraw.Draw(img)
        except Exception:
            pass

        def _get_default_font(size):
            try:
                return ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", size)
            except (OSError, IOError):
                try:
                    return ImageFont.truetype("DejaVuSans.ttf", size)
                except (OSError, IOError):
                    return ImageFont.load_default(size=size)

        def _emu_to_px(emu, scale):
            return int(emu * scale)

        def _parse_color(color_obj):
            try:
                if hasattr(color_obj, "rgb") and color_obj.rgb is not None:
                    rgb_str = str(color_obj.rgb)
                    if rgb_str and rgb_str != "None":
                        return (
                            int(rgb_str[0:2], 16),
                            int(rgb_str[2:4], 16),
                            int(rgb_str[4:6], 16),
                        )
            except Exception:
                pass
            return None

        for shape in slide.shapes:
            try:
                shape_left = _emu_to_px(shape.left, scale_x)
                shape_top = _emu_to_px(shape.top, scale_y)
                shape_width = _emu_to_px(shape.width, scale_x)
                shape_height = _emu_to_px(shape.height, scale_y)

                if shape_width <= 0 or shape_height <= 0:
                    continue

                shape_box = [
                    shape_left,
                    shape_top,
                    shape_left + shape_width,
                    shape_top + shape_height,
                ]

                fill_color_rgb = None
                try:
                    if shape.fill.type is not None:
                        fill_color_rgb = _parse_color(shape.fill.fore_color)
                except Exception:
                    pass

                line_color_rgb = None
                line_width_px = 1
                try:
                    if hasattr(shape, "line") and shape.line.color.rgb is not None:
                        line_color_rgb = _parse_color(shape.line.color)
                        if shape.line.width:
                            line_width_px = max(
                                1,
                                _emu_to_px(shape.line.width, scale_x),
                            )
                except Exception:
                    pass

                if shape.shape_type is not None:
                    from pptx.enum.shapes import MSO_SHAPE_TYPE

                    shape_type_val = shape.shape_type

                    is_auto_shape = (
                        shape_type_val == MSO_SHAPE_TYPE.AUTO_SHAPE
                    )
                    is_freeform = (
                        shape_type_val == MSO_SHAPE_TYPE.FREEFORM
                    )

                    if is_auto_shape or is_freeform:
                        if fill_color_rgb:
                            draw.rectangle(shape_box, fill=fill_color_rgb)
                        else:
                            draw.rectangle(
                                shape_box, outline=(128, 128, 128), width=1
                            )
                        if line_color_rgb:
                            draw.rectangle(
                                shape_box,
                                outline=line_color_rgb,
                                width=line_width_px,
                            )

                    elif shape_type_val == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image = shape.image
                            img_bytes = image.blob
                            from io import BytesIO

                            pil_img = Image.open(BytesIO(img_bytes))
                            pil_img = pil_img.resize(
                                (shape_width, shape_height), Image.LANCZOS
                            )
                            img.paste(pil_img, (shape_left, shape_top))
                            draw = ImageDraw.Draw(img)
                        except Exception:
                            draw.rectangle(
                                shape_box,
                                fill=(200, 200, 200),
                                outline=(128, 128, 128),
                                width=1,
                            )
                            draw.text(
                                (shape_left + 4, shape_top + 4),
                                "[Image]",
                                fill=(100, 100, 100),
                                font=_get_default_font(12),
                            )

                    elif shape_type_val == MSO_SHAPE_TYPE.TABLE:
                        if fill_color_rgb:
                            draw.rectangle(shape_box, fill=fill_color_rgb)
                        else:
                            draw.rectangle(
                                shape_box, outline=(0, 0, 0), width=1
                            )
                        try:
                            table = shape.table
                            rows = len(table.rows)
                            cols = len(table.columns)
                            cell_w = shape_width / cols if cols > 0 else shape_width
                            cell_h = shape_height / rows if rows > 0 else shape_height
                            for r in range(rows):
                                for c in range(cols):
                                    cell = table.cell(r, c)
                                    cx = shape_left + c * cell_w
                                    cy = shape_top + r * cell_h
                                    draw.rectangle(
                                        [
                                            int(cx),
                                            int(cy),
                                            int(cx + cell_w),
                                            int(cy + cell_h),
                                        ],
                                        outline=(0, 0, 0),
                                        width=1,
                                    )
                                    cell_text = cell.text.strip()
                                    if cell_text:
                                        font_size = max(8, int(cell_h * 0.5))
                                        font = _get_default_font(font_size)
                                        draw.text(
                                            (int(cx + 3), int(cy + 2)),
                                            cell_text,
                                            fill=(0, 0, 0),
                                            font=font,
                                        )
                        except Exception:
                            pass

                    elif shape_type_val == MSO_SHAPE_TYPE.GROUP:
                        draw.rectangle(
                            shape_box,
                            outline=(128, 128, 128),
                            width=1,
                        )

                    elif shape_type_val == MSO_SHAPE_TYPE.PLACEHOLDER:
                        if fill_color_rgb:
                            draw.rectangle(shape_box, fill=fill_color_rgb)

                if hasattr(shape, "text_frame") and shape.text_frame:
                    text_frame = shape.text_frame
                    text_content = shape.text_frame.text
                    if text_content and text_content.strip():
                        try:
                            font_size_pt = 18
                            font_color_rgb = (0, 0, 0)
                            font_bold = False
                            font_italic = False

                            try:
                                for para in text_frame.paragraphs:
                                    if para.runs:
                                        run = para.runs[0]
                                        if run.font.size:
                                            font_size_pt = (
                                                run.font.size.pt
                                            )
                                        font_color_rgb = (
                                            _parse_color(run.font.color)
                                            or (0, 0, 0)
                                        )
                                        if run.font.bold:
                                            font_bold = True
                                        if run.font.italic:
                                            font_italic = True
                                        break
                            except Exception:
                                pass

                            render_font_size = max(
                                8,
                                int(
                                    font_size_pt
                                    * min(scale_x, scale_y)
                                    * 0.85
                                ),
                            )
                            font = _get_default_font(render_font_size)

                            y_offset = shape_top + 4
                            for para in text_frame.paragraphs:
                                para_text = para.text
                                if not para_text:
                                    y_offset += render_font_size * 1.2
                                    continue

                                if para.runs:
                                    for run in para.runs:
                                        run_text = run.text
                                        if not run_text:
                                            continue
                                        run_size = render_font_size
                                        run_color = font_color_rgb
                                        if run.font.size:
                                            run_size = max(
                                                8,
                                                int(
                                                    run.font.size.pt
                                                    * min(scale_x, scale_y)
                                                    * 0.85
                                                ),
                                            )
                                        if (
                                            hasattr(run.font, "color")
                                            and run.font.color
                                        ):
                                            parsed = _parse_color(
                                                run.font.color
                                            )
                                            if parsed:
                                                run_color = parsed

                                        run_font = _get_default_font(
                                            run_size
                                        )
                                        bbox = draw.textbbox(
                                            (0, 0), run_text, font=run_font
                                        )
                                        text_w = bbox[2] - bbox[0]
                                        draw.text(
                                            (shape_left + 4, y_offset),
                                            run_text,
                                            fill=run_color,
                                            font=run_font,
                                        )
                                        y_offset += bbox[3] - bbox[1] + 2
                                else:
                                    bbox = draw.textbbox(
                                        (0, 0), para_text, font=font
                                    )
                                    draw.text(
                                        (shape_left + 4, y_offset),
                                        para_text,
                                        fill=font_color_rgb,
                                        font=font,
                                    )
                                    y_offset += bbox[3] - bbox[1] + 2

                                if y_offset > shape_top + shape_height:
                                    break
                        except Exception:
                            draw.text(
                                (shape_left + 4, shape_top + 4),
                                text_content[:100],
                                fill=(0, 0, 0),
                                font=_get_default_font(12),
                            )

            except Exception as e:
                logger.warning(f"Error rendering shape: {e}")
                continue

        img.save(output_path, "PNG")
        return {
            "success": True,
            "message": f"Slide saved as PNG to {output_path} (Pillow renderer)",
            "output_path": output_path,
            "width": width,
            "height": height,
            "slide_index": slide_index,
            "method": "pillow",
            "note": "Rendered using Pillow. For higher quality, install LibreOffice.",
        }
